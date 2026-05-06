"""
H-Walker 2026-05-04 First Meeting PPT Builder
사용자 16장 흐름 정확히 복원 (사용자 직접 작성한 텍스트 보존).
"""
from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from copy import deepcopy
from pathlib import Path
from lxml import etree

BASE = Path("/sessions/beautiful-dazzling-tesla/mnt/chobyeongjun--realtime-vision-control")
TEMPLATE = BASE / "docs" / "meetings" / "template.pptx"
FIG_DIR  = BASE / "docs" / "meetings" / "2026-05-04" / "figures"
OUT_PPTX = BASE / "docs" / "meetings" / "2026-05-04" / "Meeting_260504.pptx"

NS = '{http://schemas.openxmlformats.org/drawingml/2006/main}'


# ============================================================
# Helpers (Korean-safe text replacement)
# ============================================================
def set_text(shape, new_text):
    """Replace text while preserving first run's full XML (latin/ea/cs typeface, color, size)."""
    if not shape.has_text_frame:
        return False
    body = shape.text_frame._txBody
    ps = body.findall(f'{NS}p')
    if not ps:
        return False
    first_p = ps[0]
    template_pPr = first_p.find(f'{NS}pPr')
    template_r   = first_p.find(f'{NS}r')
    template_rPr = template_r.find(f'{NS}rPr') if template_r is not None else None

    for p in ps:
        body.remove(p)

    for line in new_text.split("\n"):
        p_el = etree.SubElement(body, f'{NS}p')
        if template_pPr is not None:
            p_el.append(deepcopy(template_pPr))
        r_el = etree.SubElement(p_el, f'{NS}r')
        rPr_for_line = deepcopy(template_rPr) if template_rPr is not None else etree.SubElement(
            etree.Element(f'{NS}tmp'), f'{NS}rPr'
        )
        if template_rPr is None:
            r_el.append(rPr_for_line)
        else:
            r_el.append(rPr_for_line)
        # Safety net: Korean run guarantees East-Asian typeface
        if any(0xAC00 <= ord(c) <= 0xD7A3 for c in line) and rPr_for_line.find(f'{NS}ea') is None:
            ea_el = etree.SubElement(rPr_for_line, f'{NS}ea')
            ea_el.set('typeface', 'Pretendard')
        t_el = etree.SubElement(r_el, f'{NS}t')
        t_el.text = line
    return True


def find_shape(slide, name):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    return None


def set_slide_text(slide, name, text):
    sh = find_shape(slide, name)
    if sh is None:
        print(f"  [WARN] shape '{name}' not found")
        return False
    return set_text(sh, text)


def remove_shape(shape):
    sp = shape._element
    sp.getparent().remove(sp)


def reorder_and_keep_slides(prs, keep_indices_in_order):
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    new_list = [slides_list[i] for i in keep_indices_in_order]
    for s in slides_list:
        xml_slides.remove(s)
    for s in new_list:
        xml_slides.append(s)


# ============================================================
# Build (16 slides matching user's Meeting_260504.pptx)
# ============================================================
def main():
    prs = Presentation(str(TEMPLATE))

    # ─── Slide 1 (T1 idx 0) — Cover ───
    s = prs.slides[0]
    set_slide_text(s, "MainTitle", "H-Walker\nReal-Time Vision Control")
    set_slide_text(s, "Subtitle", "케이블 드리븐 보행 보조기를 위한 실시간 하체 포즈 추정 시스템")
    set_slide_text(s, "AuthorInfo", "조병준  |  중앙대학교 기계공학과  |  AR Lab")
    set_slide_text(s, "Tag", "FIRST MEETING")

    # ─── Slide 2 (T2 idx 1) — AGENDA ───
    s = prs.slides[1]
    set_slide_text(s, "Text 12", "Introduction")
    set_slide_text(s, "Text 13", "연구 동기 및 방향")
    set_slide_text(s, "Text 16", "Experiments")
    set_slide_text(s, "Text 17", "포즈 인식 실험 – 모델, 속도, 안정성")
    set_slide_text(s, "Text 20", "Discussion")
    set_slide_text(s, "Text 21", "왜 빨라졌는가 · 실패 교훈 · 방향성")
    set_slide_text(s, "Text 24", "Next Steps")
    set_slide_text(s, "Text 25", "다음 2주 — v3 → TRT → 실모터")

    # ─── Slide 3 (T3 idx 2) — 연구 동기 ───
    # 사용자가 직접 작성한 한 단락. T3 layout 의 Text 8 (큰 본문 박스)에 모두 넣고
    # 다른 bullet 박스(9/10, 13/14, 17/18)는 비움.
    s = prs.slides[2]
    set_slide_text(s, "Text 5", "연구 동기")
    set_slide_text(s, "Text 6", "")  # 부제 자리 비움 (사용자 PPT에 없음)
    set_slide_text(s, "Text 9",  "VISION")
    set_slide_text(s, "Text 10",
        "IMU 만으로는 워커와 사람 간의 거리, 케이블 slack 정도 등 알기 힘든 값들이 많음.\n"
        "Wearable 에서는 사용할 수 없는 Walker만의 보행 상태 추정 방식\n"
        "\n"
        "→ 앞 혹은 옆에서 자세를 인식하며 실시간으로 보행 상태를 추정 가능.\n"
        "→ 비전은 다리 전체의 상태를 파악할 수 있음."
    )
    # 나머지 bullet 박스 비움
    for n in ["Text 13", "Text 14", "Text 17", "Text 18"]:
        set_slide_text(s, n, "")
    set_slide_text(s, "Text 22", "Vision")
    set_slide_text(s, "Text 23", "ZED X Mini · NVIDIA Jetson")
    set_slide_text(s, "Text 24", "Fig. 1 — Walker-specific 환경의 비전 관측")

    # ─── Slide 4 (T5 idx 4) — 연구 방향 (3축) ───
    s = prs.slides[4]
    set_slide_text(s, "Text 5", "연구 방향")
    set_slide_text(s, "Text 7", "실시간")
    set_slide_text(s, "Text 8",
        "실시간 제어를 하기위해\n"
        "사람 반응 시간 (약 100 ms) 보다 빠른 센서 입력 필요."
    )
    set_slide_text(s, "Text 9", "정확성")
    set_slide_text(s, "Text 10",
        "현재의 hip, knee, ankle joint의 위치 및 상태를 정확하게 인식"
    )
    set_slide_text(s, "Text 11", "최적의 보조")
    set_slide_text(s, "Text 12",
        "Desire Trajectory 생성 & Tracking\nForce Profile\nReinforcement Learning"
    )
    set_slide_text(s, "Text 22",
        "실시간으로 하지의 움직임을 정확하게 인지할 수 있는 포즈 인식 방식을 찾고, 이를 통해 최적의 보조를 한다."
    )
    set_slide_text(s, "Text 23", "Direction")
    set_slide_text(s, "Text 24", "Fig. 2 — 연구의 3축")

    # ─── Slide 5 (T7 idx 6) — 실험 1: 12 모델 벤치마크 표 ───
    s = prs.slides[6]
    set_slide_text(s, "Title", "실험 1 — 12개 Pose 모델 벤치마크")
    set_slide_text(s, "Header0", "지표")
    set_slide_text(s, "Header1", "YOLOv8n-Pose")
    set_slide_text(s, "Header2", "YOLO26s-Pose")
    set_slide_text(s, "Header3", "MediaPipe / RTMPose")
    set_slide_text(s, "Cell0_0", "E2E mean")
    set_slide_text(s, "Cell0_1", "약 35 ms (가장 빠름)")
    set_slide_text(s, "Cell0_2", "약 44 ms (<50 ms 아슬)")
    set_slide_text(s, "Cell0_3", "약 64 / 150+ (실패)")
    set_slide_text(s, "Cell1_0", "FPS")
    set_slide_text(s, "Cell1_1", "약 28")
    set_slide_text(s, "Cell1_2", "약 22")
    set_slide_text(s, "Cell1_3", "약 16 / 7")
    set_slide_text(s, "Cell2_0", "E2E <50 ms")
    set_slide_text(s, "Cell2_1", "전 frame 통과")
    set_slide_text(s, "Cell2_2", "거의 전 frame")
    set_slide_text(s, "Cell2_3", "사실상 실패")
    set_slide_text(s, "Cell3_0", "인식 / Conf")
    set_slide_text(s, "Cell3_1", "안정 / 0.97")
    set_slide_text(s, "Cell3_2", "안정 / 0.99")
    set_slide_text(s, "Cell3_3", "불안정 / 0.77")
    set_slide_text(s, "Footnote",
        "* Jetson Orin NX 16GB · ZED X Mini SVGA@120 · TRT FP16 · 15s 측정")

    # ─── Slide 6 (T14 idx 13) — 실험 1: E2E Latency 비교 ───
    s = prs.slides[13]
    set_slide_text(s, "Title", "실험 1 — E2E Latency 비교")
    set_slide_text(s, "Subtitle", "YOLO 계열 6종만 통과. MediaPipe · RTMPose 실패.")
    set_slide_text(s, "Caption", "Fig. 4 — 12개 Pose 모델 E2E Latency")
    imgbg = find_shape(s, "ImageBg")
    if imgbg:
        left, top, width, height = imgbg.left, imgbg.top, imgbg.width, imgbg.height
        p = find_shape(s, "ImgLabel")
        if p: remove_shape(p)
        s.shapes.add_picture(str(FIG_DIR / "fig_latency_bar.png"),
                             left, top, width=width, height=height)

    # ─── Slide 7 (T15 idx 14) — 실험 2: Fine-Tuning ───
    s = prs.slides[14]
    set_slide_text(s, "Title", "실험 2 — Fine-Tuning")
    set_slide_text(s, "Subtitle", "17 kpt → 6 kpt head 전환")
    set_slide_text(s, "Caption1", "Fig. 5 — Fine-Tuning 4대 지표 (17 → 6 kpt)")
    set_slide_text(s, "Caption2", "Fig. 6 — Evolution: 약 44 → 14 ms")
    imgbg1 = find_shape(s, "ImgBg210")
    if imgbg1:
        left, top, width, height = imgbg1.left, imgbg1.top, imgbg1.width, imgbg1.height
        p = find_shape(s, "ImgLabel_210")
        if p: remove_shape(p)
        s.shapes.add_picture(str(FIG_DIR / "fig_ft_compare.png"),
                             left, top, width=width, height=height)
    imgbg2 = find_shape(s, "ImgBg250")
    if imgbg2:
        left, top, width, height = imgbg2.left, imgbg2.top, imgbg2.width, imgbg2.height
        p = find_shape(s, "ImgLabel_250")
        if p: remove_shape(p)
        s.shapes.add_picture(str(FIG_DIR / "fig_evolution.png"),
                             left, top, width=width, height=height)

    # ─── Slide 8 (T8 idx 7) — 실험 3: Pipeline 최적화 + KPI ───
    s = prs.slides[7]
    set_slide_text(s, "Text 5", "실험 3 — Pipeline 최적화")
    set_slide_text(s, "Text 6",
        "6 kpt head 전환 + DirectTRT + C++ post → 약 44 ms 에서 약 14 ms 까지.")
    set_slide_text(s, "Text 9", "≈ 14 ms")
    set_slide_text(s, "Text 10", "처리 latency 평균")
    set_slide_text(s, "Text 11", "사람 반응의 약 1/7")
    set_slide_text(s, "Text 14", "≈ 80 Hz")
    set_slide_text(s, "Text 15", "Pipeline 처리 가능")
    set_slide_text(s, "Text 16", "publish rate에 margin")
    set_slide_text(s, "Text 19", "50–60 Hz")
    set_slide_text(s, "Text 20", "publish rate (출발선)")
    set_slide_text(s, "Text 21", "확정은 실측 후")
    set_slide_text(s, "Text 22",
        "* 180 s 연속 동작 — Hard Limit 위반 거의 없음. 정확한 publish rate은 실측 후 확정.")

    # ─── Slide 9 (T17 idx 16) — Dataset v3 (4 옷차림 카드) ───
    s = prs.slides[16]
    set_slide_text(s, "Title", "Dataset v3 — Robustness 다양성 확보")
    set_slide_text(s, "Subtitle",
        "옷차림 4종 × 보행 속도 6단계 → 1,300장 / yolo26s-lower6-v3 학습 진행")
    set_slide_text(s, "Lbl_200", "Sweats")
    set_slide_text(s, "Cap_200", "기본 보행 (2 / 3 / 4 / 4.5 / 6 / 7 km/h)")
    set_slide_text(s, "Lbl_230", "Tight")
    set_slide_text(s, "Cap_230", "신체 윤곽 강조 (3 / 4.5 / 6 km/h)")
    set_slide_text(s, "Lbl_260", "Shorts")
    set_slide_text(s, "Cap_260", "다리 노출 (3 / 4.5 / 6 km/h)")
    set_slide_text(s, "Lbl_290", "Exosuit")
    set_slide_text(s, "Cap_290", "타겟 환경 (3 / 4.5 / 6 km/h + clip)")

    # ─── Slide 10 (T11 idx 10) — 시스템 아키텍처 ───
    s = prs.slides[10]
    set_slide_text(s, "직사각형 1", "시스템 아키텍처")
    set_slide_text(s, "TextBox 2", "Perception → IPC · Safety → C++ Impedance → Motor")
    set_slide_text(s, "TextBox 9",
        "Perception 파이프라인 — Jetson 내부에서 처리되는 Pose 인식 과정")
    ph_tl = find_shape(s, "직사각형 4")
    if ph_tl:
        left, top, width, height = ph_tl.left, ph_tl.top, ph_tl.width, ph_tl.height
        for n in ["TextBox 5", "TextBox 6"]:
            p = find_shape(s, n)
            if p: remove_shape(p)
        s.shapes.add_picture(str(FIG_DIR / "fig_architecture.png"),
                             left, top, width=width, height=height)

    # ─── Slide 11 (T18 idx 17) — Evolution timeline ───
    s = prs.slides[17]
    set_slide_text(s, "직사각형 1", "Evolution")
    set_slide_text(s, "TextBox 2", "약 44 ms → 약 14 ms · Robustness 데이터 확보까지")
    set_slide_text(s, "TextBox 6",  "Phase 1")
    set_slide_text(s, "TextBox 7",  "모델 선정")
    set_slide_text(s, "TextBox 8",
        "12개 모델 벤치마크.\nYOLO 계열만\n<50 ms 통과.")
    set_slide_text(s, "TextBox 10", "Phase 2")
    set_slide_text(s, "TextBox 11", "Fine-Tuning")
    set_slide_text(s, "TextBox 12",
        "6 kpt head로 축소.\n약 44 → 18 ms.\n연산량 감소.")
    set_slide_text(s, "TextBox 14", "Phase 3")
    set_slide_text(s, "TextBox 15", "Pipeline + Stable")
    set_slide_text(s, "TextBox 16",
        "DirectTRT + C++ post\n+ patch chain.\n약 14 ms / 80 Hz.")
    set_slide_text(s, "TextBox 18", "Phase 4")
    set_slide_text(s, "TextBox 19", "Dataset v3")
    set_slide_text(s, "TextBox 20",
        "4 옷차림 × 6 속도.\n1,300장 / v3 학습.\nRobustness 확보.")

    # ─── Slide 12 (T10 idx 9) — 왜 빨라졌는가 ───
    s = prs.slides[9]
    set_slide_text(s, "직사각형 1", "왜 빨라졌는가 — 무엇을 정상화했는가")
    set_slide_text(s, "TextBox 2", "약 44 → 14 ms 가속의 원인")
    set_slide_text(s, "TextBox 7",
        "1. 연산량 감소 (-26.4 ms) : Head 17→6 kpt\n"
        "2. 파이프라인 오버랩 (-9.4 ms) : ready_rgb / depth split\n"
        "3. 프레임워크 우회 (-9 ms) : DirectTRT + C++ post\n"
        "4. GPU·OS 최대화 : jetson_clocks + MAXN + gc.off\n"
        "5. 경합 제거 : CPU isolation + SCHED_FIFO\n"
        "6. 실시간성 보장 : 20 ms HARD LIMIT + SHM seqlock"
    )
    set_slide_text(s, "TextBox 8", "— 핵심 설계 원칙")
    set_slide_text(s, "TextBox 9",
        "\"속도는 자원 분배의 문제\" — I/O는 숨기고, 프레임워크는 벗기고, 하드웨어는 고정.")

    # ─── Slide 13 (T12 idx 11) — 실패 교훈 ───
    s = prs.slides[11]
    set_slide_text(s, "직사각형 1", "실패 교훈 — 역효과였던 접근 3가지")
    set_slide_text(s, "TextBox 2", "논문 Methods 섹션의 결정적 근거로 남음")
    set_slide_text(s, "TextBox 7",  "2D 필터링 금지")
    set_slide_text(s, "TextBox 9",
        "One Euro / SegLen 을 2D에 적용 시\n"
        "depth NaN → 3D 실패.\n"
        "규칙: 필터는 3D 단계에서만.")
    set_slide_text(s, "TextBox 13", "NEURAL Depth 기각")
    set_slide_text(s, "TextBox 15",
        "2 cm 정확도 vs GPU SM 경합.\n"
        "predict 2.4× 급락, 29 FPS.\n"
        "규칙: PERFORMANCE 유지.")
    set_slide_text(s, "TextBox 19", "Zero-Copy 포기")
    set_slide_text(s, "TextBox 21",
        "copy=False → capture overwrite race.\n"
        "calib 100 → 0 %.\n"
        "규칙: 공유 버퍼는 복사 필수.")

    # ─── Slide 14 (T6 idx 5) — Next Task ───
    s = prs.slides[5]
    set_slide_text(s, "직사각형 1", "Next Task")
    set_slide_text(s, "TextBox 2",
        "다음 2주 — v3 학습 → TRT 배포 → 정확도 향상 → Teensy 데이터 전송")
    set_slide_text(s, "TextBox 7",
        "v3 학습 완료 +\nval mAP 확인")
    set_slide_text(s, "TextBox 12",
        "TRT FP16 export\n+ Pipeline 적용")
    set_slide_text(s, "TextBox 17",
        "정확도 향상")
    set_slide_text(s, "TextBox 22",
        "Teensy data 전송")
    set_slide_text(s, "TextBox 23", "Train")
    set_slide_text(s, "TextBox 24", "Deploy")
    set_slide_text(s, "TextBox 25", "Accuracy")
    set_slide_text(s, "TextBox 26", "Serial")

    # ─── Slide 15 (T13 idx 12) — 교수님 논의 사항 ───
    s = prs.slides[12]
    set_slide_text(s, "직사각형 1", "교수님 논의 사항")
    set_slide_text(s, "TextBox 2", "함께 정할 결정 3가지")
    set_slide_text(s, "TextBox 5",  "1. 제어 방식")
    set_slide_text(s, "TextBox 7",  "Desire Trajectory →\nImpedance Control")
    set_slide_text(s, "TextBox 9",  "2. 정확도 체크")
    set_slide_text(s, "TextBox 11",
        "Motion Capture / Mocap Suit으로\n"
        "정확한 위치 파악을 진행해야 하나\n"
        "or 적당히 보행을 인식할 수 있는\n"
        "정도면 괜찮을 지")
    set_slide_text(s, "TextBox 13", "3. 졸업 주제")
    set_slide_text(s, "TextBox 15",
        "이 주제로 발전시켜서\n졸업 주제로까지 이어가도\n괜찮을 지")
    set_slide_text(s, "TextBox 17", "")
    set_slide_text(s, "TextBox 19", "")

    # ─── Slide 16 (T19 idx 18) — Thanks ───
    s = prs.slides[18]
    set_slide_text(s, "Text 5",
        "AR Lab  ·  Assistive & Rehabilitation Robotics Lab  ·  Chung-Ang University  ·  조병준")

    # ============================================================
    # Reorder & Delete unused (16 slides)
    # ============================================================
    # 매핑 (사용자 16장 → template index):
    #  1 Cover (0)
    #  2 AGENDA (1)
    #  3 연구 동기 — Vision/IMU 어필 (2)         T3
    #  4 연구 방향 — 3축 (4)                       T5
    #  5 실험 1 표 (6)                             T7
    #  6 실험 1 latency bar (13)                   T14
    #  7 실험 2 — Fine-Tuning (14)                 T15
    #  8 실험 3 — Pipeline + KPI (7)              T8 3 KPI
    #  9 Dataset v3 — 4 옷차림 (16)                T17
    # 10 시스템 아키텍처 (10)                      T11
    # 11 Evolution timeline (17)                   T18
    # 12 왜 빨라졌는가 (9)                          T10
    # 13 실패 교훈 (11)                             T12
    # 14 Next Task (5)                             T6 process flow
    # 15 교수님 논의 사항 (12)                      T13 4-feature
    # 16 Thanks (18)                                T19
    # 미사용: 3 (T4 two-column), 8 (T9 4 KPI), 15 (T16 3 photos)

    desired = [0, 1, 2, 4, 6, 13, 14, 7, 16, 10, 17, 9, 11, 5, 12, 18]
    reorder_and_keep_slides(prs, desired)

    OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPTX))
    print(f"Saved: {OUT_PPTX}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
